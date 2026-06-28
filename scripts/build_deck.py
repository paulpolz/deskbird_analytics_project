#!/usr/bin/env python3
"""Build deck.pptx from deck.md with editable text and movable PNG charts.

Marp --pptx rasterizes each slide; this script keeps text editable and embeds
chart PNGs as separate shapes (from export_slide_assets.py).

Usage:
  python scripts/build_deck.py
  make deck
"""

from __future__ import annotations

import re
import sys
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
SLIDES_DIR = ROOT / "docs" / "slides"
DECK_MD = SLIDES_DIR / "deck.md"
OUT_PPTX = SLIDES_DIR / "deck.pptx"
ASSETS = SLIDES_DIR / "assets"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MARGIN_IN = 0.45
CONTENT_W_IN = SLIDE_W_IN - 2 * MARGIN_IN
CANVAS_PX = 1280
FONT_FAMILY = "Berlin Type"

IMG_RE = re.compile(r"!\[(?:width:(\d+)px)?\]\(([^)]+)\)")


@dataclass
class SlideContent:
  title: str = ""
  subtitle: str = ""
  blocks: list[tuple[str, object]] = field(default_factory=list)


def _strip_md_inline(text: str) -> str:
  text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
  text = re.sub(r"\*(.+?)\*", r"\1", text)
  text = re.sub(r"`(.+?)`", r"\1", text)
  return text.strip()


def _px_to_inches(px: int) -> float:
  return (px / CANVAS_PX) * SLIDE_W_IN


def _parse_table(lines: list[str]) -> list[list[str]] | None:
  rows = [ln for ln in lines if ln.strip().startswith("|")]
  if len(rows) < 2:
    return None
  parsed: list[list[str]] = []
  for row in rows:
    if re.match(r"^\|[\s\-:|]+\|$", row.strip()):
      continue
    cells = [_strip_md_inline(c.strip()) for c in row.strip().strip("|").split("|")]
    parsed.append(cells)
  return parsed or None


def _parse_slide(raw: str) -> SlideContent:
  slide = SlideContent()
  lines = raw.splitlines()
  text_buf: list[str] = []

  def flush_text() -> None:
    if not text_buf:
      return
    table = _parse_table(text_buf)
    if table:
      prose = [ln for ln in text_buf if not ln.strip().startswith("|")]
      if prose:
        slide.blocks.append(("text", prose))
      slide.blocks.append(("table", table))
    else:
      slide.blocks.append(("text", text_buf.copy()))
    text_buf.clear()

  for line in lines:
    img_matches = list(IMG_RE.finditer(line))
    if img_matches:
      flush_text()
      images: list[tuple[Path, float]] = []
      for m in img_matches:
        px = int(m.group(1) or 900)
        rel = m.group(2)
        images.append((SLIDES_DIR / rel, _px_to_inches(px)))
      slide.blocks.append(("images", images))
      continue

    if line.startswith("# "):
      slide.title = _strip_md_inline(line[2:])
    elif line.startswith("## "):
      slide.subtitle = _strip_md_inline(line[3:])
    elif line.strip():
      text_buf.append(line.rstrip())

  flush_text()
  return slide


def _load_slides(path: Path) -> list[SlideContent]:
  text = path.read_text(encoding="utf-8")
  if text.startswith("---"):
    end = text.find("---", 3)
    if end != -1:
      text = text[end + 3 :].lstrip("\n")
  chunks = [c.strip() for c in text.split("\n---\n") if c.strip()]
  return [_parse_slide(c) for c in chunks]


def _set_font(para, *, size: int, bold: bool = False) -> None:
  para.font.name = FONT_FAMILY
  para.font.size = Pt(size)
  para.font.bold = bold
  para.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)


def _add_textbox(slide, left, top, width, height, text: str, *, font_size=14, bold=False):
  box = slide.shapes.add_textbox(left, top, width, height)
  tf = box.text_frame
  tf.word_wrap = True
  tf.vertical_anchor = MSO_ANCHOR.TOP
  p = tf.paragraphs[0]
  p.text = text
  _set_font(p, size=font_size, bold=bold)
  return box


def _add_body_textbox(slide, left, top, width, height, lines: list[str]):
  box = slide.shapes.add_textbox(left, top, width, height)
  tf = box.text_frame
  tf.word_wrap = True
  tf.vertical_anchor = MSO_ANCHOR.TOP

  first = True
  for raw in lines:
    line = raw.rstrip()
    if not line.strip():
      continue
    bullet = line.lstrip().startswith("- ")
    text = _strip_md_inline(line.lstrip("- ").strip())
    if not text:
      continue

    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.text = text
    p.level = 0 if bullet else 0
    _set_font(p, size=14)
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    if bullet:
      p.level = 0
    if text.startswith("Why it matters:") or text.startswith("Actions:"):
      p.font.bold = True
  return box


def _add_table(slide, left, top, width, rows: list[list[str]]):
  if not rows:
    return
  n_rows, n_cols = len(rows), max(len(r) for r in rows)
  height = Inches(0.35 * n_rows + 0.2)
  shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
  table = shape.table

  for r, row in enumerate(rows):
    for c in range(n_cols):
      cell = table.cell(r, c)
      cell.text = row[c] if c < len(row) else ""
      for para in cell.text_frame.paragraphs:
        _set_font(para, size=12 if r else 13, bold=r == 0)
  return shape


def _layout_images(slide, images: list[tuple[Path, float]], top_in: float) -> float:
  if not images:
    return top_in

  gap_in = 0.25
  n = len(images)

  if n == 1:
    path, w_in = images[0]
    w = min(w_in, CONTENT_W_IN)
    left = Inches(MARGIN_IN + (CONTENT_W_IN - w) / 2)
    if path.exists():
      pic = slide.shapes.add_picture(str(path), left, Inches(top_in), width=Inches(w))
      return top_in + pic.height.inches + gap_in
    return top_in + 3.2 + gap_in

  slot_w = (CONTENT_W_IN - gap_in * (n - 1)) / n
  x = MARGIN_IN + max(0.0, (CONTENT_W_IN - n * slot_w - gap_in * (n - 1)) / 2)
  bottom = top_in

  for path, w_in in images:
    w = min(w_in, slot_w)
    if path.exists():
      pic = slide.shapes.add_picture(str(path), Inches(x), Inches(top_in), width=Inches(w))
      bottom = max(bottom, top_in + pic.height.inches)
    else:
      print(f"  warning: missing {path.relative_to(ROOT)}", file=sys.stderr)
    x += w + gap_in

  return bottom + gap_in


def _estimate_text_height(lines: list[str]) -> float:
  chars = sum(len(_strip_md_inline(ln.lstrip("- ").strip())) for ln in lines if ln.strip())
  lines_est = max(1, chars // 90 + sum(1 for ln in lines if ln.lstrip().startswith("- ")))
  return 0.22 * lines_est + 0.15


def _build_slide(prs: Presentation, content: SlideContent) -> None:
  layout = prs.slide_layouts[6]
  slide = prs.slides.add_slide(layout)
  y = MARGIN_IN

  if content.title:
    _add_textbox(slide, Inches(MARGIN_IN), Inches(y), Inches(CONTENT_W_IN), Inches(0.55), content.title, font_size=28, bold=True)
    y += 0.65

  if content.subtitle:
    subtitle_h = min(1.4, max(0.45, 0.08 * len(content.subtitle) / 10 + 0.35))
    _add_textbox(slide, Inches(MARGIN_IN), Inches(y), Inches(CONTENT_W_IN), Inches(subtitle_h), content.subtitle, font_size=16)
    y += subtitle_h + 0.15

  for kind, payload in content.blocks:
    if kind == "text":
      lines = payload  # type: ignore[assignment]
      h = _estimate_text_height(lines)
      remaining = SLIDE_H_IN - y - MARGIN_IN
      if remaining > 0.25:
        _add_body_textbox(slide, Inches(MARGIN_IN), Inches(y), Inches(CONTENT_W_IN), Inches(min(h, remaining)), lines)
        y += min(h, remaining) + 0.1
    elif kind == "images":
      y = _layout_images(slide, payload, y)  # type: ignore[arg-type]
    elif kind == "table":
      rows = payload  # type: ignore[assignment]
      _add_table(slide, Inches(MARGIN_IN), Inches(y), Inches(CONTENT_W_IN), rows)
      y += 0.35 * len(rows) + 0.35


def build_deck(md_path: Path = DECK_MD, out_path: Path = OUT_PPTX) -> Path:
  slides = _load_slides(md_path)
  prs = Presentation()
  prs.slide_width = Inches(SLIDE_W_IN)
  prs.slide_height = Inches(SLIDE_H_IN)

  for content in slides:
    _build_slide(prs, content)

  prs.save(out_path)
  return out_path


def main() -> None:
  if not DECK_MD.exists():
    print(f"Missing {DECK_MD}", file=sys.stderr)
    sys.exit(1)

  out = build_deck()
  rel = out.relative_to(ROOT)
  print(f"Wrote {rel} ({len(_load_slides(DECK_MD))} slides, editable text + PNG shapes)")


if __name__ == "__main__":
  main()
